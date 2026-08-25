from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

OUT = "Delta_AI_Newsletter_專案介紹.pptx"
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

FONT = "Microsoft JhengHei"
C = {
    "dark": "071C2C", "navy": "0B3954", "ink": "102A43", "teal": "00A896",
    "cyan": "8DE3E0", "mint": "D8F3DC", "orange": "FFB703", "coral": "F07167",
    "cream": "F7F4ED", "white": "FFFFFF", "muted": "637381", "line": "D7E2EA",
    "pale": "EAF4F4", "panel": "12334A",
}

def rgb(hexstr):
    return RGBColor.from_string(hexstr)

def fill(shape, color, transparency=0):
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(color)
    if transparency: shape.fill.transparency = transparency
    shape.line.color.rgb = rgb(color)

def text(slide, value, x, y, w, h, size=18, color="ink", bold=False, align=PP_ALIGN.LEFT,
         valign=MSO_ANCHOR.MIDDLE, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True; tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0; tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = value
    run.font.name = FONT; run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = rgb(C.get(color, color))
    return box

def rect(slide, x, y, w, h, color="white", line=None, radius=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    fill(shape, C.get(color, color))
    if line:
        shape.line.color.rgb = rgb(C.get(line, line)); shape.line.width = Pt(0.8)
    return shape

def line(slide, x1, y1, x2, y2, color="teal", width=1.2):
    shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    shape.line.color.rgb = rgb(C.get(color, color)); shape.line.width = Pt(width)
    return shape

def pill(slide, value, x, y, w, color="teal", bg="pale"):
    rect(slide, x, y, w, 0.34, bg, bg, True)
    text(slide, value, x+0.08, y+0.04, w-0.16, 0.2, 9, color, True, PP_ALIGN.CENTER)

def card(slide, x, y, w, h, color="white", border="line"):
    return rect(slide, x, y, w, h, color, border, True)

def header(slide, title, dark=False, section="DELTA AI NEWSLETTER"):
    text(slide, section, 0.62, 0.3, 5, 0.2, 9, "cyan" if dark else "teal", True)
    text(slide, title, 0.62, 0.68, 12, 0.55, 27, "white" if dark else "ink", True)

def footer(slide, n, dark=False):
    line(slide, 0.62, 6.98, 12.72, 6.98, "panel" if dark else "line", 0.7)
    text(slide, f"DELTA AI NEWSLETTER  /  {n:02d}", 0.62, 7.08, 4, 0.18, 8, "muted" if not dark else "cyan", True)

def bg(slide, color="cream"):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(C[color])

def add_slide(color="cream"):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, color); return s

# 1 Cover
s = add_slide("dark")
rect(s, 8.95, 0, 4.38, 7.5, "teal")
line(s, 9.25, 1.15, 12.2, 4.65, "cyan", 1.1); line(s, 10.05, 4.3, 12.4, 1.25, "cyan", 1.1)
for x, y, r, c in [(9.3,1.2,.18,"orange"),(11.0,2.5,.13,"white"),(10.05,4.22,.18,"cyan"),(12.05,4.85,.11,"orange")]:
    sh=s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(r), Inches(r)); fill(sh,C[c])
text(s, "DELTA AI NEWSLETTER", .72, .68, 5, .3, 11, "cyan", True)
text(s, "把外部 AI 情報\n變成可讀、可追溯的內容", .72, 1.7, 7.7, 1.45, 34, "white", True, valign=MSO_ANCHOR.TOP)
text(s, "專案全貌與目前進度", .76, 3.6, 4, .35, 18, "mint")
text(s, "產品／專案團隊介紹｜2026.08", .76, 6.35, 4.5, .25, 11, "cyan")
text(s, "01", 11.8, 6.32, .7, .35, 22, "white", True, PP_ALIGN.RIGHT)

# 2 Problem
s = add_slide(); header(s, "為什麼需要這套系統")
text(s, "不是缺資訊，而是缺少「跟我有關」的判斷。", .62, 1.5, 8.8, .5, 24, "ink", True)
text(s, "台達超過 25 個單位關注點不同；只做熱門度排行，內容很容易變成大家都看過、但沒有人覺得與工作相關。", .62, 2.18, 7.1, .65, 16, "muted", valign=MSO_ANCHOR.TOP)
for i, (a,b,c) in enumerate([("資訊太多","每天持續產生新內容","coral"),("需求分散","能源、資安、研發各看不同角度","orange"),("品質風險","AI 生成前必須先限制素材邊界","teal")]):
    sh=s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(.72), Inches(3.35+i*.85), Inches(.18), Inches(.18)); fill(sh,C[c])
    text(s,a,1.08,3.22+i*.85,1.25,.3,14,"ink",True); text(s,b,2.48,3.22+i*.85,4.6,.3,13,"muted")
card(s,8.55,1.3,3.95,4.6,"dark","dark")
text(s,"INPUT",8.92,1.7,1.3,.2,10,"cyan",True); text(s,"18 個公開來源",8.92,2.15,2.8,.45,22,"white",True)
text(s,"每天新增約 33 篇\n最多一天 74 篇",8.92,3.08,2.9,.6,15,"mint",valign=MSO_ANCHOR.TOP)
text(s,"如果沒有篩選，\n讀者只會得到更多噪音。",9.0,4.75,2.95,.7,15,"white",True,valign=MSO_ANCHOR.TOP)
footer(s,2)

# 3 Goal
s = add_slide("dark"); header(s, "系統目標", True, "THE PRODUCT IDEA")
text(s,"重點在準，不在多。",.62,1.55,7.2,.6,34,"white",True)
text(s,"每天從公開來源抓取 AI 內容，先判斷、再聚類、再選題，最後寫成中文短文與英文版。",.66,2.38,6.3,.7,16,"cyan",valign=MSO_ANCHOR.TOP)
for i,(num,label) in enumerate([("9","期日報已產出"),("26","篇文章已生成"),("18","個讀者模組"),("25+","個台達關注單位")]):
    x=.66+i*2.05; text(s,num,x,4.15,1.7,.55,31,"orange" if i==1 else "cyan",True); text(s,label,x,4.9,1.7,.3,12,"cyan")
card(s,8.1,1.35,4.25,4.7,"panel","panel")
text(s,"一句話理解",8.55,1.82,2.5,.25,11,"orange",True); text(s,"把「外部情報」\n變成「內部可用的判斷」",8.55,2.5,3.25,1.05,25,"white",True,valign=MSO_ANCHOR.TOP)
text(s,"不是單純摘要機，而是一條有門檻、有配額、有紀錄的內容生產線。",8.55,4.45,3.15,.7,15,"mint",valign=MSO_ANCHOR.TOP)
footer(s,3,True)

# 4 Lifecycle
s=add_slide(); header(s,"每天怎麼跑")
text(s,"從來源到出刊，分成「建池」與「出刊」兩段。",.62,1.42,8.5,.45,22,"ink",True)
pill(s,"建池｜每日排程",.68,2.2,1.65,"navy","pale")
steps=[("01","抓取","18 個來源"),("02","關卡","日期／長度"),("03","聚類","同一件事合併"),("04","標籤","摘要／關鍵詞"),("05","打分","18 個模組")]
for i,(n,a,b) in enumerate(steps):
    x=.68+i*2.38; card(s,x,2.82,1.95,1.3,"pale" if i==2 else "white")
    text(s,n,x+.17,3.04,.35,.22,12,"teal",True); text(s,a,x+.58,3.0,1.1,.25,15,"ink",True); text(s,b,x+.17,3.5,1.6,.2,10,"muted")
    if i<4: text(s,"›",x+2.05,3.22,.2,.3,24,"orange",True,PP_ALIGN.CENTER)
pill(s,"出刊｜需要時執行",.68,4.55,1.75,"navy","pale")
steps=[("06","選題","輪動＋配額"),("07","取素材","只用話題自己的文章"),("08","寫作","中文短文＋洞察"),("09","自檢","不過就重寫／不出"),("10","翻譯","英文版預先完成")]
for i,(n,a,b) in enumerate(steps):
    x=.68+i*2.38; card(s,x,5.15,1.95,1.3,"FFF4D6" if i==3 else "white")
    text(s,n,x+.17,5.37,.35,.22,12,"coral" if i==3 else "teal",True); text(s,a,x+.58,5.33,1.1,.25,15,"ink",True); text(s,b,x+.17,5.83,1.62,.3,10,"muted",valign=MSO_ANCHOR.TOP)
    if i<4: text(s,"›",x+2.05,5.55,.2,.3,24,"orange",True,PP_ALIGN.CENTER)
footer(s,4)

# 5 Topic
s=add_slide(); header(s,"核心設計")
text(s,"處理單位是「話題」，不是「文章」。",.62,1.46,8.9,.5,25,"ink",True)
text(s,"同一件事的多篇報導先聚在一起，再一起打分、選題與寫作；避免重複，也讓素材邊界清楚。",.62,2.12,7.2,.6,15,"muted",valign=MSO_ANCHOR.TOP)
card(s,.72,3.22,3.05,2.38); text(s,"原始輸入",1.05,3.5,1.3,.22,11,"muted",True)
for i,(k,v) in enumerate([("A","TechCrunch 報導"),("B","企業案例文章"),("C","產業媒體訊號")]):
    y=3.98+i*.47; rect(s,1.02,y,2.45,.33,"pale","line",True); text(s,k,1.18,y+.06,.2,.16,9,"teal",True); text(s,v,1.48,y+.05,1.7,.18,11,"ink")
text(s,"→",4.15,4.25,.55,.35,28,"orange",True,PP_ALIGN.CENTER)
card(s,5.18,3.02,3,2.78,"dark","dark"); text(s,"TOPIC",5.55,3.38,1.1,.22,10,"cyan",True); text(s,"同一件事\n聚成一個話題",5.55,3.85,2.3,.8,23,"white",True,valign=MSO_ANCHOR.TOP); text(s,"embedding + Qdrant\n本機完成語意比對",5.55,5.05,2.2,.45,12,"mint",valign=MSO_ANCHOR.TOP)
text(s,"→",8.6,4.25,.55,.35,28,"orange",True,PP_ALIGN.CENTER)
card(s,9.6,3.22,2.95,2.38); text(s,"可用輸出",9.93,3.5,1.3,.22,11,"muted",True); text(s,"一個清楚的主軸",9.93,4.05,2.2,.32,18,"navy",True); text(s,"話題自己的來源\n→ 寫作素材邊界",9.93,4.72,2.15,.5,13,"muted",valign=MSO_ANCHOR.TOP)
text(s,"目前 326 個話題中，94.8% 底下只有一篇文章，這是現況，不是隱藏的缺口。",.72,6.18,11.5,.25,11,"muted",italic=True)
footer(s,5)

# 6 Modules
s=add_slide("dark"); header(s,"選題邏輯",True,"THE EDITORIAL ENGINE")
text(s,"18 個模組，代表 18 種讀者視角。",.62,1.46,8.8,.5,26,"white",True); text(s,"不是單純挑最高分，而是讓不同業務與職能都得到版位，同時保留判斷的可追溯性。",.62,2.12,7.25,.55,15,"cyan",valign=MSO_ANCHOR.TOP)
text(s,"職能視角  10",.72,3.08,2.1,.22,11,"orange",True)
for i,v in enumerate(["法務合規","財會稽核","人資","行銷品牌","資訊資安","研發工程","營運物流","策略投資","知識管理","EHS"]): pill(s,v,.72+(i%5)*1.48,3.5+(i//5)*.52,1.27,"mint","panel")
text(s,"本業視角  8",.72,4.72,2.1,.22,11,"cyan",True)
for i,v in enumerate(["能源電力","樓宇自動化","電動車車用","網通基礎設施","製造廠務","消費性產品","軟體平台","永續節能"]): pill(s,v,.72+(i%4)*1.78,5.15+(i//4)*.52,1.58,"cyan","panel")
card(s,9.15,2.9,3.25,3.2,"panel","panel"); text(s,"三輪選題",9.55,3.32,1.5,.25,12,"orange",True)
for i,(n,v) in enumerate([("01","模組輪流推薦"),("02","總分高者遞補"),("03","版位不足才保底")]): text(s,f"{n}   {v}",9.55,3.88+i*.68,2.5,.25,14,"white" if i else "cyan",True)
text(s,"本業模組保底過半，避免整期變成通用工具介紹。",9.55,5.28,2.25,.45,12,"mint",valign=MSO_ANCHOR.TOP)
footer(s,6,True)

# 7 Gates
s=add_slide(); header(s,"品質控制")
text(s,"三道 Gate，讓「為什麼沒上刊」有答案。",.62,1.46,8.9,.5,25,"ink",True); text(s,"落選不是資料消失，而是留下狀態與理由；「不夠格」和「版位滿」分開統計。",.62,2.12,7.5,.55,15,"muted",valign=MSO_ANCHOR.TOP)
gates=[("GATE 1","文章能不能用","過舊、太短、與 AI 無關","included / signal_only / excluded","teal"),("GATE 2","話題有沒有可寫內容","沒有實質內文就不進候選","單一來源保留，但如實揭露","orange"),("GATE 3","這篇能不能出刊","分數、素材、自檢、版位","理由寫入 selection_trace","coral")]
for i,(a,b,c,d,col) in enumerate(gates):
    x=.72+i*4.1; card(s,x,3.24,3.55,2.62,"white",col); rect(s,x,3.24,.1,2.62,col,col); text(s,a,x+.35,3.62,1.1,.2,10,col,True); text(s,b,x+.35,4.08,2.6,.35,18,"ink",True); text(s,c,x+.35,4.7,2.7,.3,13,"muted"); text(s,d,x+.35,5.25,2.75,.35,11,"ink",valign=MSO_ANCHOR.TOP)
text(s,"每個候選話題都有去向：入選記錄模組、分數、輪次；落選記錄理由。",.72,6.3,11.5,.25,12,"navy",True); footer(s,7)

# 8 Architecture
s=add_slide(); header(s,"系統架構")
text(s,"LLM 負責判斷與生成，本機元件負責相似度與資料保存。",.62,1.46,10,.5,23,"ink",True)
rows=[("來源",[("公開網站","18 sources"),("RSS／API","RawItem")]),("處理",[("Python pipeline","ingestion / pipeline"),("LLM gateway","tag / score / write")]),("儲存",[("SQLite","topics.db"),("Qdrant","embedding vectors"),("Neo4j（選配）","knowledge graph")]),("呈現",[("本機網頁","issues / trace"),("英文版","預先翻譯完成")])]
for ri,(label,items) in enumerate(rows):
    y=2.25+ri*.98; text(s,label,.72,y+.23,1.1,.22,11,"teal",True)
    for i,(a,b) in enumerate(items):
        x=2+i*3.35; card(s,x,y,2.75,.75); text(s,a,x+.22,y+.14,2.25,.2,14,"ink",True); text(s,b,x+.22,y+.43,2.25,.16,10,"muted")
text(s,"資料不離開公司：LLM 走台達內網 gateway；embedding、聚類、reranker 在本機。",.72,6.35,11.6,.25,12,"navy",True); footer(s,8)

# 9 Progress / lesson
s=add_slide("dark"); header(s,"目前進度與重要學習",True,"WHAT WE HAVE LEARNED")
text(s,"系統已經跑通，但品質還需要人與評測來閉環。",.62,1.46,10.7,.5,24,"white",True)
for i,(num,label) in enumerate([("345","篇文章進池"),("326","個話題"),("9","期日報"),("2,463","筆 LLM 呼叫紀錄")]):
    x=.72+(i%2)*2.35; y=2.55+(i//2)*1.35; text(s,num,x,y,1.7,.5,29,"orange" if i==2 else "cyan",True); text(s,label,x,y+.62,1.8,.22,12,"cyan")
card(s,5.75,2.35,6.55,3.5,"panel","panel"); text(s,"最重要的品質發現",6.2,2.77,2.9,.25,12,"orange",True)
text(s,"過去的素材補全會把\n「主題相關」誤當成「同一件事」。",6.2,3.35,5,.9,23,"white",True,valign=MSO_ANCHOR.TOP)
text(s,"修正：寫作素材現在只用話題自己的文章。\n舊有 26 篇文章中，20 篇需要下架重寫。",6.2,4.72,5.2,.65,14,"mint",valign=MSO_ANCHOR.TOP)
footer(s,9,True)

# 10 Next steps
s=add_slide(); header(s,"下一步與會議重點")
text(s,"下一階段不是再加功能，而是先把品質與治理補完整。",.62,1.46,10.7,.5,24,"ink",True)
items=[("01","人工審核流程","從單一旗標升級成核准／退回，錯的內容不能直接出門。","coral"),("02","建立評測集","用 before / after 數字驗證打分、AI 判定與聚類門檻。","orange"),("03","重跑與重寫","修正早期標籤；26 篇已出刊內容中 20 篇需重寫後再審。","teal"),("04","權限與刊期","選題紀錄頁加登入；機器每日跑，刊物規劃回到週報。","navy")]
for i,(n,a,b,col) in enumerate(items):
    x=.72+(i%2)*6.1; y=2.35+(i//2)*1.75; card(s,x,y,5.45,1.35,"white",col); text(s,n,x+.32,y+.28,.45,.25,13,col,True); text(s,a,x+.95,y+.22,3.8,.28,17,"ink",True); text(s,b,x+.95,y+.65,4.1,.42,12,"muted",valign=MSO_ANCHOR.TOP)
rect(s,.72,6.12,11.55,.52,"dark","dark",True); text(s,"今天要對齊的結論：先完成人工審核與評測，再決定是否擴大發信與解凍知識圖譜。",.98,6.25,11,.22,14,"white",True,PP_ALIGN.CENTER)
footer(s,10)

prs.save(OUT)
print(OUT)
